from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = Path(r"C:\코딩\로고\전환설계연구소_PPT-템플릿_v0.5.pptx")
OUT = Path(r"C:\코딩\교육설계\output\slides\CRAFTO_RADAR_AGENTSMD_설명실습_전용템플릿.pptx")

NAVY = RGBColor(0x1C, 0x2F, 0x5C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x61, 0x61, 0x61)
MID_GREY = RGBColor(0x94, 0x94, 0x94)
LIGHT_GREY = RGBColor(0xF4, 0xF5, 0xF7)
LINE = RGBColor(0xDA, 0xDA, 0xDA)
LIGHT_BLUE = RGBColor(0xDC, 0xE4, 0xF0)
PALE_BLUE = RGBColor(0xF0, 0xF4, 0xFA)
GREEN = RGBColor(0x5C, 0x7A, 0x69)
AMBER = RGBColor(0xB3, 0x7B, 0x3F)

FONT = "Pretendard"
SERIF = "Playfair Display"


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=24,
    color=BLACK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_multiline(
    slide,
    lines,
    x,
    y,
    w,
    h,
    size=20,
    color=BLACK,
    bold_first=False,
    gap_pt=8,
    font=FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bool(bold_first and idx == 0)
        p.space_after = Pt(gap_pt)
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False, width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(width)
    return shp


def add_line(slide, x, y, w, color=LINE, weight=1.0):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_brand(slide, page=None, section="CRAFTO-RADAR"):
    add_text(slide, "전환설계연구소", 0.72, 0.34, 2.3, 0.25, 10, NAVY, bold=True)
    add_text(slide, section, 9.45, 0.34, 2.45, 0.25, 8.5, GREY, font=SERIF, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.72, 0.70, 11.9, color=LINE)
    if page:
        add_text(slide, str(page).zfill(2), 12.18, 6.82, 0.32, 0.16, 7.5, GREY, font=SERIF, align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_text(slide, kicker, 0.72, 0.95, 4.0, 0.24, 10, NAVY, bold=True, font=SERIF)
    add_text(slide, title, 0.72, 1.27, 10.25, 0.68, 30, NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.97, 10.8, 0.38, 15.5, GREY)


def add_pill(slide, text, x, y, w, fill=NAVY, color=WHITE, size=10.5):
    shp = add_rect(slide, x, y, w, 0.32, fill=fill, line=fill, radius=True)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = True
    return shp


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_text(slide, "전환설계연구소", 0.72, 0.55, 2.5, 0.3, 11, NAVY, bold=True)
    add_text(slide, "AGENTIC AI WORKSHOP MODULE", 0.72, 1.42, 4.2, 0.26, 10.5, GREY, font=SERIF)
    add_text(slide, "CRAFTO-RADAR", 0.70, 2.05, 8.9, 0.8, 42, NAVY, bold=True, font=SERIF)
    add_text(slide, "LLM 질문 설계에서 에이전트 작업공간 설계로", 0.74, 2.92, 8.6, 0.45, 18, BLACK)
    add_text(slide, "강사용 설명 · 실습 슬라이드", 0.75, 5.63, 4.2, 0.28, 12, GREY)
    add_text(slide, "2026.05", 10.35, 5.63, 1.9, 0.28, 12, GREY, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.72, 4.55, 5.75, color=NAVY, weight=2)
    add_rect(slide, 9.0, 1.05, 2.85, 4.35, fill=PALE_BLUE, line=PALE_BLUE)
    add_text(slide, "C", 9.45, 1.47, 0.5, 0.38, 21, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "맥락", 10.05, 1.53, 1.05, 0.24, 11.5, GREY)
    add_line(slide, 9.35, 2.15, 1.92, color=LINE)
    add_text(slide, "R", 9.45, 2.55, 0.5, 0.38, 21, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "범위", 10.05, 2.61, 1.05, 0.24, 11.5, GREY)
    add_line(slide, 9.35, 3.23, 1.92, color=LINE)
    add_text(slide, "A", 9.45, 3.63, 0.5, 0.38, 21, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "승인", 10.05, 3.69, 1.05, 0.24, 11.5, GREY)
    add_line(slide, 9.35, 4.31, 1.92, color=LINE)
    add_text(slide, "D", 9.45, 4.71, 0.5, 0.38, 21, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "완료", 10.05, 4.77, 1.05, 0.24, 11.5, GREY)


def slide_shift(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 2)
    add_title(slide, "질문을 잘 쓰는 일에서 일을 맡기는 일로", "Agentic AI에서는 답변 품질보다 작업 범위와 승인 흐름이 더 중요해집니다.", "WHY")
    add_text(slide, "LLM 활용", 1.05, 2.85, 3.5, 0.36, 22, NAVY, bold=True)
    add_multiline(slide, ["좋은 답을 받는다", "맥락을 정리한다", "초안과 피드백을 주고받는다"], 1.05, 3.42, 3.95, 1.2, 17, BLACK, gap_pt=9)
    add_text(slide, "→", 6.05, 3.54, 0.6, 0.5, 30, MID_GREY, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "Agentic AI 활용", 7.25, 2.85, 4.2, 0.36, 22, NAVY, bold=True)
    add_multiline(slide, ["파일을 만들고 고친다", "폴더와 자료를 다룬다", "계획, 승인, 보고가 필요하다"], 7.25, 3.42, 4.25, 1.2, 17, BLACK, gap_pt=9)
    add_text(slide, "핵심 변화", 1.08, 5.72, 1.2, 0.25, 11, GREY, font=SERIF)
    add_text(slide, "프롬프트 역량이 작업공간 설계 역량으로 확장됩니다.", 2.28, 5.66, 8.8, 0.35, 18, NAVY, bold=True)


def slide_crafto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 3)
    add_title(slide, "CRAFTO는 그대로 유지합니다", "프롬프트 공식이 아니라, 내 생각을 AI에게 전달하는 맥락 구조입니다.", "REVIEW")
    items = [
        ("C", "Context", "상황과 목적"),
        ("R", "Role", "AI의 역할"),
        ("A", "Audience", "결과물을 볼 사람"),
        ("F", "Format", "형식과 구조"),
        ("T", "Tone", "문체와 톤"),
        ("O", "Option", "제약과 옵션"),
    ]
    for i, (letter, term, desc) in enumerate(items):
        x = 0.95 + (i % 3) * 3.9
        y = 2.70 + (i // 3) * 1.48
        add_text(slide, letter, x, y, 0.55, 0.42, 24, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
        add_text(slide, term, x + 0.72, y + 0.02, 2.1, 0.28, 15.5, BLACK, bold=True)
        add_text(slide, desc, x + 0.72, y + 0.42, 2.6, 0.28, 13.5, GREY)
        add_line(slide, x, y + 0.88, 2.95, color=LINE)
    add_text(slide, "CRAFTO는 좋은 답을 받기 위한 출발점입니다.", 2.1, 5.92, 8.9, 0.32, 17, NAVY, bold=True, align=PP_ALIGN.CENTER)


def slide_radar_need(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 4)
    add_title(slide, "Agentic AI에는 RADAR가 추가됩니다", "AI가 실제 파일을 다룰 때는 맥락만으로 충분하지 않습니다.", "EXTENSION")
    add_text(slide, "CRAFTO", 1.15, 2.85, 3.1, 0.48, 25, NAVY, bold=True, font=SERIF)
    add_text(slide, "맥락 설계", 1.18, 3.42, 2.6, 0.32, 16, BLACK, bold=True)
    add_text(slide, "무엇을 왜 요청하는지 정리합니다.", 1.18, 3.82, 3.45, 0.5, 14.5, GREY)
    add_text(slide, "+", 5.65, 3.22, 0.45, 0.42, 24, MID_GREY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "RADAR", 7.0, 2.85, 3.1, 0.48, 25, NAVY, bold=True, font=SERIF)
    add_text(slide, "실행 설계", 7.03, 3.42, 2.6, 0.32, 16, BLACK, bold=True)
    add_text(slide, "어디까지 맡기고 어떻게 점검할지 정합니다.", 7.03, 3.82, 3.85, 0.5, 14.5, GREY)
    add_line(slide, 1.1, 5.28, 10.85, color=NAVY, weight=2)
    add_text(slide, "CRAFTO-RADAR는 에이전트에게 일을 맡기는 요청서입니다.", 1.45, 5.62, 10.1, 0.34, 18, NAVY, bold=True, align=PP_ALIGN.CENTER)


def slide_radar_terms(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 5)
    add_title(slide, "RADAR는 AGENTS.md를 만드는 질문 프레임입니다", "작업 범위, 승인, 완료, 질문, 보고를 정하면 에이전트가 안정적으로 움직입니다.", "FRAME")
    items = [
        ("R", "Range", "어디까지 읽고 쓸 수 있는가"),
        ("A", "Approval", "무엇은 승인 후 해야 하는가"),
        ("D", "Done", "무엇이 나오면 완료인가"),
        ("A", "Ask", "언제 멈추고 질문해야 하는가"),
        ("R", "Report", "작업 후 무엇을 보고해야 하는가"),
    ]
    start_x = 0.95
    for i, (letter, term, desc) in enumerate(items):
        x = start_x + i * 2.35
        add_text(slide, letter, x, 2.72, 0.55, 0.5, 30, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.1, 3.34, 0.35, color=NAVY, weight=2)
        add_text(slide, term, x - 0.12, 3.66, 1.25, 0.26, 12.5, BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x - 0.34, 4.10, 1.7, 0.76, 12.2, GREY, align=PP_ALIGN.CENTER)
    add_text(slide, "짧게 말하면", 1.35, 5.70, 1.35, 0.22, 10.5, GREY, font=SERIF)
    add_text(slide, "AI에게 줄 작업실 사용 규칙을 다섯 문장으로 정리하는 일입니다.", 2.68, 5.62, 8.85, 0.34, 18, NAVY, bold=True)


def slide_agents_md(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 6)
    add_title(slide, "AGENTS.md는 폴더 안의 운영 문서입니다", "프롬프트가 이번 한 번의 요청이라면, AGENTS.md는 이 폴더에서 계속 적용되는 작업 규칙입니다.", "WORKSPACE")
    add_rect(slide, 1.05, 2.68, 4.75, 2.42, fill=LIGHT_GREY, line=LINE)
    add_text(slide, "프롬프트", 1.38, 3.02, 3.8, 0.36, 22, NAVY, bold=True)
    add_text(slide, "이번 요청의 맥락", 1.38, 3.55, 3.5, 0.3, 16, BLACK, bold=True)
    add_text(slide, "대화창 안에서 주고받는 지시입니다.", 1.38, 3.95, 3.85, 0.45, 14, GREY)
    add_rect(slide, 7.0, 2.68, 4.75, 2.42, fill=PALE_BLUE, line=LIGHT_BLUE)
    add_text(slide, "AGENTS.md", 7.33, 3.02, 3.8, 0.36, 22, NAVY, bold=True, font=SERIF)
    add_text(slide, "작업공간의 규칙", 7.33, 3.55, 3.5, 0.3, 16, BLACK, bold=True)
    add_text(slide, "폴더 안에서 에이전트가 따르는 운영 문서입니다.", 7.33, 3.95, 3.85, 0.45, 14, GREY)
    add_text(slide, "AGENTS.md는 RADAR 질문에 답하면서 만듭니다.", 2.42, 5.86, 8.5, 0.32, 17, NAVY, bold=True, align=PP_ALIGN.CENTER)


def slide_practice_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 7)
    add_title(slide, "홈페이지 실습은 이 순서로 진행합니다", "바로 만들지 않고, 먼저 작업공간을 설계합니다.", "PRACTICE")
    steps = [
        ("01", "CRAFTO", "홈페이지 요청의 맥락을 정리"),
        ("02", "RADAR", "AGENTS.md 작성 조건 정리"),
        ("03", "폴더 지정", "input · wiki · output 준비"),
        ("04", "계획 승인", "작업계획을 받고 승인"),
        ("05", "파일 제작", "index.html · style.css 생성"),
        ("06", "검토", "자료 밖 표현과 수정 항목 점검"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.92 + (i % 3) * 3.85
        y = 2.62 + (i // 3) * 1.55
        add_text(slide, num, x, y, 0.48, 0.26, 13, NAVY, bold=True, font=SERIF)
        add_text(slide, head, x + 0.58, y, 2.55, 0.30, 17, BLACK, bold=True)
        add_text(slide, body, x + 0.58, y + 0.42, 2.75, 0.42, 12.8, GREY)
        add_line(slide, x, y + 1.03, 3.1, color=LINE)
    add_text(slide, "성능 체감은 결과물보다 ‘계획 → 승인 → 실행 → 검토’ 흐름에서 생깁니다.", 1.58, 6.0, 9.7, 0.3, 16, NAVY, bold=True, align=PP_ALIGN.CENTER)


def slide_agents_template(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 8)
    add_title(slide, "RADAR로 AGENTS.md 초안을 만듭니다", "참가자는 빈칸을 자기 강의·프로젝트에 맞게 채웁니다.", "TEMPLATE")
    rows = [
        ("Range", "input 자료만 읽고 output 폴더에 결과물을 만든다."),
        ("Approval", "파일 생성·수정 전 작업계획을 먼저 제시한다."),
        ("Done", "index.html, style.css, 검토메모.md를 만든다."),
        ("Ask", "자료에 없는 성과·가격·후기는 질문한다."),
        ("Report", "생성 파일, 주요 구성, 최종 결정 항목을 보고한다."),
    ]
    for i, (head, body) in enumerate(rows):
        y = 2.55 + i * 0.66
        add_text(slide, head, 1.0, y, 1.45, 0.25, 14, NAVY, bold=True, font=SERIF)
        add_text(slide, body, 2.62, y, 8.8, 0.28, 14.8, BLACK)
        add_line(slide, 1.0, y + 0.43, 10.65, color=LINE)
    add_text(slide, "강의 포인트", 1.05, 6.0, 1.3, 0.24, 10.5, GREY, font=SERIF)
    add_text(slide, "잘 만든 AGENTS.md는 긴 프롬프트를 줄이고, 반복 작업의 품질을 안정화합니다.", 2.35, 5.93, 9.0, 0.32, 16.5, NAVY, bold=True)


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 9)
    add_text(slide, "정리", 0.72, 1.05, 1.2, 0.28, 10.5, NAVY, bold=True, font=SERIF)
    add_text(slide, "CRAFTO는 질문의 맥락을 정리합니다.", 1.05, 2.00, 10.8, 0.55, 30, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "RADAR는 에이전트의 실행 조건을 정리합니다.", 1.05, 2.86, 10.8, 0.55, 30, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "AGENTS.md는 그 내용을 폴더 안에 남기는 운영 문서입니다.", 2.0, 4.05, 9.15, 0.38, 18, BLACK, align=PP_ALIGN.CENTER)
    add_line(slide, 3.25, 5.18, 6.8, color=NAVY, weight=2)
    add_text(slide, "AI에게 일을 맡긴다는 것은 요청을 길게 쓰는 일이 아니라, 일할 조건을 설계하는 일입니다.", 1.65, 5.58, 10.0, 0.34, 16.5, GREY, align=PP_ALIGN.CENTER)


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    slide_cover(prs)
    slide_shift(prs)
    slide_crafto(prs)
    slide_radar_need(prs)
    slide_radar_terms(prs)
    slide_agents_md(prs)
    slide_practice_flow(prs)
    slide_agents_template(prs)
    slide_close(prs)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
