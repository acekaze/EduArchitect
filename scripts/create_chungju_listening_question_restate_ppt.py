from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


TEMPLATE = Path(r"C:\코딩\로고\전환설계연구소_PPT-템플릿_v0.5.pptx")
OUT = Path(r"C:\코딩\교육설계\output\slides\충주_퍼실리테이터_경청질문재진술_전용템플릿.pptx")

NAVY = RGBColor(0x1C, 0x2F, 0x5C)
GREY = RGBColor(0x61, 0x61, 0x61)
LIGHT_BLUE = RGBColor(0xDC, 0xE4, 0xF0)
LINE = RGBColor(0xDA, 0xDA, 0xDA)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF7, 0xF8, 0xFA)
GREEN = RGBColor(0x5C, 0x7A, 0x69)
AMBER = RGBColor(0xB3, 0x7B, 0x3F)

FONT = "Pretendard"
SERIF = "Playfair Display"


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=24,
    color=BLACK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_multiline(
    slide,
    lines,
    x,
    y,
    w,
    h,
    size=22,
    color=BLACK,
    bold_first=False,
    bullet=False,
    gap_pt=8,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bool(bold_first and idx == 0)
        p.space_after = Pt(gap_pt)
        if bullet:
            p.level = 0
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def add_brand(slide, page=None, section="CHUNGJU FACILITATOR"):
    add_text(slide, "전환설계연구소", 0.72, 0.34, 2.2, 0.28, 10, NAVY, bold=True)
    add_text(slide, section, 9.65, 0.34, 2.2, 0.28, 8.5, GREY, font=SERIF, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(0.70), Inches(11.9), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE
    if page:
        add_text(slide, str(page).zfill(2), 12.18, 6.82, 0.32, 0.16, 7.5, GREY, font=SERIF, align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_text(slide, kicker, 0.72, 0.95, 3.5, 0.28, 10, NAVY, bold=True, font=SERIF)
    add_text(slide, title, 0.72, 1.27, 9.6, 0.65, 30, NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.95, 9.7, 0.42, 16, GREY)


def add_pill(slide, text, x, y, w, fill=NAVY, color=WHITE):
    shp = add_rect(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    shp.text_frame.clear()
    p = shp.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.color.rgb = color
    run.font.bold = True
    return shp


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_text(slide, "전환설계연구소", 0.72, 0.55, 2.5, 0.3, 11, NAVY, bold=True)
    add_text(slide, "WORKSHOP MODULE", 0.72, 1.55, 3.1, 0.26, 11, GREY, font=SERIF)
    add_text(slide, "경청 · 질문 · 재진술", 0.70, 2.05, 9.8, 0.8, 40, NAVY, bold=True)
    add_text(slide, "존중을 행동으로 바꾸는 대화 기본기", 0.74, 2.92, 7.8, 0.45, 18, BLACK)
    add_text(slide, "충주 퍼실리테이터 과정", 0.75, 5.63, 4.0, 0.28, 12, GREY)
    add_text(slide, "2026.05.13", 10.35, 5.63, 1.9, 0.28, 12, GREY, align=PP_ALIGN.RIGHT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(4.55), Inches(5.8), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.color.rgb = NAVY
    soft = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.08), Inches(3.25), Inches(4.25))
    soft.fill.solid()
    soft.fill.fore_color.rgb = LIGHT_BLUE
    soft.line.color.rgb = LIGHT_BLUE
    add_text(slide, "LISTEN", 9.12, 1.45, 2.6, 0.5, 22, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "ASK", 9.12, 2.62, 2.6, 0.5, 22, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(slide, "RESTATE", 9.12, 3.79, 2.6, 0.5, 22, NAVY, bold=True, font=SERIF, align=PP_ALIGN.CENTER)


def slide_statement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 2)
    add_text(slide, "존중은 태도에서 끝나지 않습니다.", 1.05, 1.75, 10.8, 0.62, 34, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "상대가 말할 시간을 주고,\n들은 내용을 함께 다룰 수 있게 만드는 행동입니다.",
        2.05,
        2.72,
        8.7,
        1.18,
        24,
        BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "오늘은 그 행동을 세 가지 기술로 연습합니다.", 3.25, 4.68, 6.5, 0.35, 16, GREY, align=PP_ALIGN.CENTER)
    for i, word in enumerate(["경청", "질문", "재진술"]):
        x = 3.25 + i * 2.2
        add_pill(slide, word, x, 5.23, 1.38, fill=NAVY)


def slide_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 3)
    add_title(slide, "대화가 업무로 이어지는 순서", "많이 듣고, 제대로 나누고, 빈칸을 묻고, 함께 다룰 문장으로 다시 말합니다.")
    steps = [
        ("01", "많이 듣기", "상대가 말할 시간을 확보한다"),
        ("02", "제대로 듣기", "감정·사실·요구·조치로 나눈다"),
        ("03", "질문하기", "빠진 정보를 채운다"),
        ("04", "재진술하기", "함께 다룰 문장으로 바꾼다"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = 0.92 + i * 3.05
        add_text(slide, num, x, 3.05, 0.7, 0.34, 16, NAVY, bold=True, font=SERIF)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.52), Inches(2.25), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = NAVY if i < 2 else GREY
        line.line.color.rgb = NAVY if i < 2 else GREY
        add_text(slide, title, x, 3.82, 2.2, 0.36, 18, BLACK, bold=True)
        add_text(slide, desc, x, 4.24, 2.25, 0.72, 13.5, GREY)
        if i < len(steps) - 1:
            add_text(slide, "→", x + 2.42, 3.84, 0.3, 0.32, 18, GREY, align=PP_ALIGN.CENTER)


def slide_two_listenings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 4)
    add_title(slide, "경청은 두 가지로 나눕니다", "많이 듣는 태도와 제대로 나누는 기술이 함께 있어야 합니다.")
    add_text(slide, "많이 듣기", 1.0, 2.78, 4.8, 0.45, 24, NAVY, bold=True)
    add_multiline(
        slide,
        ["끊지 않는다", "반응을 짧게 준다", "침묵을 급하게 채우지 않는다", "상대가 끝까지 말하게 둔다"],
        1.0,
        3.40,
        4.55,
        1.75,
        18,
        BLACK,
        gap_pt=9,
    )
    add_text(slide, "제대로 듣기", 7.0, 2.78, 4.8, 0.45, 24, NAVY, bold=True)
    add_multiline(
        slide,
        ["감정과 사실을 나눈다", "요구와 조치를 구분한다", "빠진 정보를 찾는다", "마지막 문장으로 정리한다"],
        7.0,
        3.40,
        4.55,
        1.75,
        18,
        BLACK,
        gap_pt=9,
    )
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.35), Inches(2.75), Inches(0.02), Inches(2.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LINE
    bar.line.color.rgb = LINE
    add_text(slide, "태도", 1.0, 5.74, 1.2, 0.26, 11, GREY, font=SERIF)
    add_text(slide, "구조", 7.0, 5.74, 1.2, 0.26, 11, GREY, font=SERIF)


def slide_classification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 5)
    add_title(slide, "제대로 듣기는 네 칸으로 나누는 일입니다", "상대 말이 복잡할수록 한 덩어리로 받지 않습니다.")
    items = [
        ("감정", "불편함, 서운함, 답답함"),
        ("사실", "상대가 겪었다고 말하는 일"),
        ("요구", "바뀌길 원하는 내용"),
        ("조치", "실제로 처리할 다음 행동"),
    ]
    for i, (head, body) in enumerate(items):
        x = 0.85 + (i % 2) * 5.95
        y = 2.72 + (i // 2) * 1.45
        add_rect(slide, x, y, 5.15, 0.95, fill=SOFT, line=LINE)
        add_text(slide, head, x + 0.25, y + 0.18, 1.35, 0.38, 20, NAVY, bold=True)
        add_text(slide, body, x + 1.68, y + 0.24, 3.0, 0.3, 15, BLACK)
    add_text(slide, "분류가 끝나면 바로 답하지 말고, 비어 있는 칸을 먼저 찾습니다.", 1.55, 5.55, 9.6, 0.35, 17, GREY, align=PP_ALIGN.CENTER)


def slide_practice_sentence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 6)
    add_title(slide, "실습 | 이 말을 네 칸으로 나눕니다", "먼저 개인으로 작성하고, 이후 조별로 서로의 분류를 비교합니다.")
    add_text(
        slide,
        "지난번 프로그램 안내가 너무 늦었어요.\n현장에 갔더니 담당자마다 말이 달라서\n기다리는 사람들만 불편했습니다.\n다음에는 이런 식으로 운영하면 안 됩니다.",
        1.1,
        2.68,
        10.6,
        2.0,
        25,
        BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "작성 순서: 감정 → 사실 → 요구 → 조치", 3.5, 5.42, 5.7, 0.35, 16, NAVY, bold=True, align=PP_ALIGN.CENTER)


def add_example_item(slide, num, title, quote, x, y, w=10.85):
    add_text(slide, num, x, y, 0.45, 0.25, 13, NAVY, bold=True, font=SERIF)
    add_text(slide, title, x + 0.55, y, 2.6, 0.28, 15.5, NAVY, bold=True)
    add_text(slide, quote, x + 3.15, y - 0.02, w - 3.15, 0.58, 12.7, BLACK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.72), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def slide_civic_examples_a(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 7)
    add_title(slide, "문화재단형 민원 사례 6개 | 1", "각 조는 하나를 골라 감정·사실·요구·조치로 나눕니다.")
    examples = [
        (
            "01",
            "안내 부족",
            "신청은 했는데 안내 문자가 너무 늦게 왔어요. 장소도 정확히 모르겠고 주차 안내도 없어서 한참 헤맸습니다. 다음에는 미리 제대로 안내해 주세요.",
        ),
        (
            "02",
            "현장 대기",
            "예약하고 왔는데도 현장에서 너무 오래 기다렸습니다. 접수 줄도 길고 안내하는 사람도 없어서 왜 예약을 받았는지 모르겠더라고요.",
        ),
        (
            "03",
            "참여 기회",
            "이런 프로그램은 항상 아는 사람들만 먼저 신청하는 것 같아요. 공지가 어디에 올라오는지도 모르겠고, 일반 시민은 참여하기가 어렵습니다.",
        ),
    ]
    for i, item in enumerate(examples):
        add_example_item(slide, *item, x=0.95, y=2.72 + i * 1.12)
    add_text(slide, "선택 후 먼저 감정과 사실을 분리합니다.", 3.35, 6.10, 6.2, 0.26, 14.5, GREY, align=PP_ALIGN.CENTER)


def slide_civic_examples_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 8)
    add_title(slide, "문화재단형 민원 사례 6개 | 2", "분류가 끝나면 빠진 정보를 질문으로 채웁니다.")
    examples = [
        (
            "04",
            "난이도 차이",
            "초보자도 참여할 수 있다고 해서 왔는데 설명이 너무 어려웠습니다. 아이와 같이 왔는데 따라가기 힘들었어요.",
        ),
        (
            "05",
            "응대 혼선",
            "문의했을 때 담당자마다 답변이 달랐습니다. 어떤 분은 된다고 하고, 다른 분은 안 된다고 하니 시민 입장에서는 너무 혼란스럽습니다.",
        ),
        (
            "06",
            "결과 공유",
            "시민 의견을 받는다고 해서 설문도 했는데 이후에 뭐가 반영됐는지 알 수가 없습니다. 의견만 받고 끝나는 느낌입니다.",
        ),
    ]
    for i, item in enumerate(examples):
        add_example_item(slide, *item, x=0.95, y=2.72 + i * 1.12)
    add_text(slide, "마지막에는 재진술 1문장과 다음 조치 1문장을 만듭니다.", 2.55, 6.10, 7.8, 0.26, 14.5, GREY, align=PP_ALIGN.CENTER)


def slide_classification_answer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 9)
    add_title(slide, "분류 예시", "정답을 맞히는 활동이 아니라, 대화에서 무엇을 다뤄야 하는지 찾는 활동입니다.")
    rows = [
        ("감정", "안내 지연과 현장 혼선에 대한 불편"),
        ("사실", "안내가 늦었고 담당자마다 설명이 달랐다고 말함"),
        ("요구", "다음 회차에서는 같은 문제가 없길 원함"),
        ("조치", "안내 방식, 현장 동선, 담당자 응대를 조정해야 함"),
    ]
    for i, (head, body) in enumerate(rows):
        y = 2.58 + i * 0.78
        add_text(slide, head, 1.02, y, 1.25, 0.3, 17, NAVY, bold=True)
        add_text(slide, body, 2.4, y, 8.7, 0.34, 17, BLACK)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y + 0.43), Inches(10.8), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.color.rgb = LINE
    add_text(slide, "여기까지가 경청입니다. 다음은 비어 있는 정보를 질문으로 채웁니다.", 1.4, 5.85, 9.8, 0.35, 16, GREY, align=PP_ALIGN.CENTER)


def slide_find_gaps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 10)
    add_title(slide, "분류 후에는 빈칸을 찾습니다", "무엇이 부족한지 알아야 질문이 짧고 정확해집니다.")
    rows = [
        ("감정은 강한데", "상황이 불분명하다", "언제, 어디서 있었나요?"),
        ("사실은 있는데", "영향이 빠져 있다", "어떤 점이 가장 불편하셨나요?"),
        ("요구는 있는데", "가능한 조치가 모호하다", "원하시는 다음 조치는 무엇인가요?"),
        ("조치는 있는데", "담당과 일정이 없다", "누가 언제까지 안내하면 좋을까요?"),
    ]
    headers = ["들은 내용", "비어 있는 정보", "질문으로 전환"]
    xs = [0.95, 4.10, 7.25]
    ws = [2.6, 2.65, 4.05]
    for x, w, h in zip(xs, ws, headers):
        add_text(slide, h, x, 2.44, w, 0.3, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = 2.93 + i * 0.72
        for j, text in enumerate(row):
            add_rect(slide, xs[j], y, ws[j], 0.48, fill=SOFT if j < 2 else LIGHT_BLUE, line=WHITE)
            add_text(slide, text, xs[j] + 0.1, y + 0.12, ws[j] - 0.2, 0.18, 12.4, BLACK if j < 2 else NAVY, bold=(j == 2), align=PP_ALIGN.CENTER)


def slide_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 11)
    add_title(slide, "질문은 추측을 정보로 바꾸는 기술입니다", "상대를 캐묻는 것이 아니라, 함께 다룰 수 있는 말로 바꾸는 과정입니다.")
    items = [
        ("상황", "그 일이 언제, 어디서 있었나요?"),
        ("영향", "그 일 때문에 어떤 점이 가장 불편하셨나요?"),
        ("요구", "원하시는 조치는 무엇인가요?"),
        ("실행", "다음에는 어떤 방식으로 안내받으면 좋으실까요?"),
    ]
    for i, (head, question) in enumerate(items):
        y = 2.58 + i * 0.84
        add_text(slide, head, 1.1, y, 1.05, 0.28, 16, NAVY, bold=True)
        add_text(slide, question, 2.35, y, 8.7, 0.33, 20, BLACK)
    add_text(slide, "좋은 질문은 길지 않습니다. 빠진 정보 한 가지를 정확히 묻습니다.", 1.45, 5.85, 9.7, 0.3, 15.5, GREY, align=PP_ALIGN.CENTER)


def slide_question_rules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 12)
    add_title(slide, "질문할 때 피해야 할 것과 바꿀 것", "질문은 상대를 압박하지 않고 대화를 앞으로 움직여야 합니다.")
    pairs = [
        ("왜 그렇게 하셨어요?", "그때 어떤 상황이었나요?"),
        ("그건 오해 아닌가요?", "어느 부분에서 그렇게 느끼셨나요?"),
        ("원래 그런 건데요.", "어떤 안내가 있었으면 좋으셨나요?"),
    ]
    add_text(slide, "피해야 할 질문", 1.1, 2.48, 4.6, 0.34, 17, AMBER, bold=True)
    add_text(slide, "바꿔 말하기", 7.0, 2.48, 4.6, 0.34, 17, GREEN, bold=True)
    for i, (bad, better) in enumerate(pairs):
        y = 3.05 + i * 0.78
        add_text(slide, bad, 1.1, y, 4.7, 0.34, 18, BLACK)
        add_text(slide, "→", 6.1, y, 0.4, 0.34, 18, GREY, align=PP_ALIGN.CENTER)
        add_text(slide, better, 7.0, y, 4.8, 0.34, 18, NAVY, bold=True)
    add_text(slide, "핵심은 책임 추궁보다 상황 파악입니다.", 3.4, 5.82, 6.2, 0.3, 16, GREY, align=PP_ALIGN.CENTER)


def slide_restate(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 13)
    add_title(slide, "재진술은 반복이 아니라 정리입니다", "상대 말을 그대로 따라 하는 것이 아니라, 함께 다룰 수 있는 문장으로 바꿉니다.")
    add_rect(slide, 1.05, 2.65, 10.8, 1.35, fill=SOFT, line=LINE)
    add_text(
        slide,
        "말씀을 정리하면,\n______ 때문에 불편하셨고,\n특히 ______ 부분에서 어려움이 있었으며,\n앞으로는 ______을 원하신다는 뜻으로 이해했습니다.",
        1.45,
        2.88,
        10.0,
        0.95,
        20,
        BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "재진술 후에는 반드시 상대에게 맞는지 묻습니다.", 3.15, 4.88, 6.6, 0.35, 17, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "\"제가 이해한 내용이 맞을까요?\"", 4.0, 5.36, 4.9, 0.34, 19, GREY, align=PP_ALIGN.CENTER)


def slide_restate_practice(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide)
    add_brand(slide, 14)
    add_title(slide, "실습 | 듣고, 묻고, 다시 말합니다", "마지막 산출물은 재진술 1문장과 다음 조치 1문장입니다.")
    add_text(slide, "1", 1.0, 2.75, 0.4, 0.3, 18, NAVY, bold=True, font=SERIF)
    add_text(slide, "상대 발화를 네 칸으로 나눈다", 1.48, 2.75, 4.8, 0.3, 18, BLACK, bold=True)
    add_text(slide, "2", 1.0, 3.45, 0.4, 0.3, 18, NAVY, bold=True, font=SERIF)
    add_text(slide, "빠진 정보를 질문 3개로 채운다", 1.48, 3.45, 4.8, 0.3, 18, BLACK, bold=True)
    add_text(slide, "3", 1.0, 4.15, 0.4, 0.3, 18, NAVY, bold=True, font=SERIF)
    add_text(slide, "재진술 1문장으로 정리한다", 1.48, 4.15, 4.8, 0.3, 18, BLACK, bold=True)
    add_text(slide, "4", 1.0, 4.85, 0.4, 0.3, 18, NAVY, bold=True, font=SERIF)
    add_text(slide, "다음 조치 1문장으로 마무리한다", 1.48, 4.85, 4.8, 0.3, 18, BLACK, bold=True)
    add_rect(slide, 7.2, 2.62, 4.55, 2.5, fill=LIGHT_BLUE, line=LIGHT_BLUE)
    add_text(slide, "마무리 문장", 7.55, 2.92, 3.9, 0.3, 16, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "오늘 정리한 내용은\n담당자와 공유하고,\n______까지 ______ 방식으로\n안내드리겠습니다.", 7.58, 3.42, 3.85, 1.15, 18, NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "들은 말이 다음 행동으로 이어질 때 대화가 일합니다.", 2.25, 5.9, 8.5, 0.3, 17, GREY, align=PP_ALIGN.CENTER)


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    slide_cover(prs)
    slide_statement(prs)
    slide_flow(prs)
    slide_two_listenings(prs)
    slide_classification(prs)
    slide_practice_sentence(prs)
    slide_civic_examples_a(prs)
    slide_civic_examples_b(prs)
    slide_classification_answer(prs)
    slide_find_gaps(prs)
    slide_questions(prs)
    slide_question_rules(prs)
    slide_restate(prs)
    slide_restate_practice(prs)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
